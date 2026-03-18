import json
import csv
import io
from pathlib import Path

from app.core.logging import log
from app.utils.filetypes import get_extension


def load_text_from_bytes(content: bytes, filename: str) -> str:
    """Extract text from file bytes based on file extension."""
    ext = get_extension(filename)
    text = ""

    try:
        if ext == ".pdf":
            text = _load_pdf(content)
        elif ext == ".csv":
            text = _load_csv(content)
        elif ext == ".json":
            text = _load_json(content)
        elif ext in (".doc", ".docx"):
            text = _load_docx(content)
        elif ext == ".pptx":
            text = _load_pptx(content)
        elif ext in (".xlsx", ".xls"):
            text = _load_xlsx(content)
        elif ext == ".rtf":
            text = _load_rtf(content)
        elif ext in (".html", ".htm"):
            text = _load_html(content)
        elif ext == ".xml":
            text = _load_xml(content)
        else:
            text = content.decode("utf-8", errors="replace")
    except Exception as e:
        log.warning(f"Failed to parse {filename} as {ext}, falling back to raw decode: {e}")
        text = content.decode("utf-8", errors="replace")

    return text


def load_text_from_file(path: Path) -> str:
    """Load text from a file on disk."""
    content = path.read_bytes()
    return load_text_from_bytes(content, path.name)


def _load_pdf(content: bytes) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)
    except ImportError:
        log.warning("PyPDF2 not installed — decoding PDF as raw text")
        return content.decode("utf-8", errors="replace")


def _load_docx(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))

        return "\n\n".join(paragraphs)
    except ImportError:
        log.warning("python-docx not installed — cannot read .docx files")
        return content.decode("utf-8", errors="replace")


def _load_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text)
    except ImportError:
        log.warning("python-pptx not installed — cannot read .pptx files")
        return content.decode("utf-8", errors="replace")


def _load_xlsx(content: bytes) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets_text)
    except ImportError:
        log.warning("openpyxl not installed — cannot read .xlsx files")
        return content.decode("utf-8", errors="replace")


def _load_rtf(content: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        rtf_content = content.decode("utf-8", errors="replace")
        return rtf_to_text(rtf_content)
    except ImportError:
        log.warning("striprtf not installed — cannot read .rtf files")
        return content.decode("utf-8", errors="replace")


def _load_html(content: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        text = content.decode("utf-8", errors="replace")
        soup = BeautifulSoup(text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        log.warning("beautifulsoup4 not installed — stripping tags manually")
        import re
        text = content.decode("utf-8", errors="replace")
        return re.sub(r'<[^>]+>', '', text)


def _load_xml(content: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        text = content.decode("utf-8", errors="replace")
        soup = BeautifulSoup(text, "html.parser")
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        import re
        text = content.decode("utf-8", errors="replace")
        return re.sub(r'<[^>]+>', '', text)


def _load_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        rows.append(" | ".join(row))
    return "\n".join(rows)


def _load_json(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    try:
        data = json.loads(text)
        return _json_to_text(data)
    except json.JSONDecodeError:
        return text


def _flatten_value(value, depth: int = 0) -> str:
    """Recursively convert a JSON value into a readable inline string."""
    if isinstance(value, dict):
        pairs = [f"{k}: {_flatten_value(v, depth + 1)}" for k, v in value.items()]
        sep = "; " if depth > 0 else "\n"
        return sep.join(pairs)
    if isinstance(value, list):
        if not value:
            return ""
        # Homogeneous list of scalars → comma-joined
        if all(isinstance(i, (str, int, float, bool)) or i is None for i in value):
            return ", ".join("null" if i is None else str(i) for i in value)
        # Heterogeneous or nested → one entry per line
        return "\n".join(_flatten_value(i, depth + 1) for i in value)
    if value is None:
        return "null"
    return str(value)


def _json_to_text(data) -> str:
    """Convert parsed JSON into chunking-friendly plain text.

    Strategies:
    - List of objects  → one paragraph per object (natural chunk boundaries).
    - Dict             → key: value pairs, nested arrays of objects expanded.
    - Scalar / other   → plain string conversion.
    """
    if isinstance(data, list):
        blocks: list[str] = []
        for item in data:
            if isinstance(item, dict):
                lines = [f"{k}: {_flatten_value(v, depth=1)}" for k, v in item.items()]
                blocks.append("\n".join(lines))
            else:
                blocks.append(_flatten_value(item))
        return "\n\n".join(blocks)

    if isinstance(data, dict):
        parts: list[str] = []
        for k, v in data.items():
            if isinstance(v, list) and v and all(isinstance(i, dict) for i in v):
                # Nested array of objects — give each its own labelled block
                sub_blocks = [f"[{k}]"]
                for item in v:
                    lines = [f"  {ik}: {_flatten_value(iv, depth=1)}" for ik, iv in item.items()]
                    sub_blocks.append("\n".join(lines))
                parts.append("\n".join(sub_blocks))
            elif isinstance(v, dict):
                lines = [f"  {sk}: {_flatten_value(sv, depth=1)}" for sk, sv in v.items()]
                parts.append(f"{k}:\n" + "\n".join(lines))
            else:
                parts.append(f"{k}: {_flatten_value(v, depth=1)}")
        return "\n\n".join(parts)

    return str(data)
